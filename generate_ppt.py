import os
import subprocess
import sys

def install_and_import(package):
    try:
        __import__(package)
    except ImportError:
        print(f"Installing {package}...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", package])
        
install_and_import('pptx')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define brand colors
    NAVY = RGBColor(8, 28, 56)
    GOLD = RGBColor(212, 175, 55)
    WHITE = RGBColor(255, 255, 255)

    # Helper function to apply dark background to a slide
    def apply_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Layout
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "مشروع المنصة الرقمية الموحدة"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    title.text_frame.paragraphs[0].font.name = 'Arial'
    
    subtitle.text = "معهد إعداد القادة بحلوان\nوزارة التعليم العالي والبحث العلمي - 2026"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'

    # --- Slide 2: Project Objectives ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide)
    
    title = slide.shapes.title
    title.text = "أهداف التطوير ورؤية المنصة الجديدة"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    
    content = slide.placeholders[1]
    content.text_frame.clear()  # Clear placeholder
    
    p = content.text_frame.add_paragraph()
    p.text = "1. تصميم رئاسي فاخر يعكس هيبة وعراقة المعهد منذ 1919."
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    p2 = content.text_frame.add_paragraph()
    p2.text = "2. تجربة مستخدم (UI/UX) عالمية بتأثيرات بصرية متقدمة."
    p2.font.color.rgb = WHITE
    p2.space_after = Pt(20)
    
    p3 = content.text_frame.add_paragraph()
    p3.text = "3. استمارة تسجيل ذكية لدعم قواعد بيانات الطلاب بدقة."
    p3.font.color.rgb = WHITE
    p3.space_after = Pt(20)

    p4 = content.text_frame.add_paragraph()
    p4.text = "4. عرض متجاوب 100% للشاشات والهواتف الذكية لدعم شريحة الشباب."
    p4.font.color.rgb = WHITE

    # --- Artifact Paths for Images ---
    # We will use the screenshots captured by the browser subagent
    artifact_dir = r"C:\Users\Khale\.gemini\antigravity\brain\15216448-59aa-4f61-b7e7-c8762f51fed2"
    hero_img = os.path.join(artifact_dir, "hero_section_view_1773997915260.png")
    about_img = os.path.join(artifact_dir, "about_section_view_1773997955983.png")
    programs_img = os.path.join(artifact_dir, "programs_section_view_1773998006747.png")

    # --- Slide 3: Hero Section Preview ---
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide)
    
    title = slide.shapes.title
    title.text = "الواجهة الرئيسية (البانر الذهبي) - Hero Section"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    
    if os.path.exists(hero_img):
        # Add image to slide, centering it
        slide.shapes.add_picture(hero_img, Inches(0.5), Inches(2), width=Inches(9))

    # --- Slide 4: About Section Preview ---
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide)
    
    title = slide.shapes.title
    title.text = "تاريخ المعهد والأرقام (الرؤية والرسالة) - About Section"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    
    if os.path.exists(about_img):
        slide.shapes.add_picture(about_img, Inches(0.5), Inches(2), width=Inches(9))

    # --- Slide 5: Programs Section Preview ---
    slide = prs.slides.add_slide(slide_layout)
    apply_dark_bg(slide)
    
    title = slide.shapes.title
    title.text = "المسارات الاستراتيجية (البرامج التدريبية)"
    title.text_frame.paragraphs[0].font.color.rgb = GOLD
    
    if os.path.exists(programs_img):
        slide.shapes.add_picture(programs_img, Inches(0.5), Inches(2), width=Inches(9))

    # Save presentation
    output_filename = r"c:\Users\Khale\OneDrive\Desktop\معهد اعداد القاده\Presentation_معهد_اعداد_القادة.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved to {output_filename}")

if __name__ == "__main__":
    create_presentation()
